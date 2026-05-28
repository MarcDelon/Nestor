import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfdoc
from reportlab.pdfgen import canvas

pdfdoc.md5 = lambda *args, **kwargs: hashlib.md5(*args)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs"
OUT.mkdir(exist_ok=True)

PPTX_PATH = OUT / "SafeTrip_Audit_Plan_Amelioration.pptx"
PDF_PATH = OUT / "SafeTrip_Checklist_Technique.pdf"

CREAM = RGBColor(0xF7, 0xF4, 0xEE)
CREAM_2 = RGBColor(0xED, 0xE9, 0xDF)
FOREST = RGBColor(0x0A, 0x2F, 0x1D)
EMERALD = RGBColor(0x00, 0x67, 0x3C)
YELLOW = RGBColor(0xFC, 0xD1, 0x16)
RED = RGBColor(0xCE, 0x11, 0x26)
INK = RGBColor(0x1C, 0x2B, 0x22)
MUTED = RGBColor(0x71, 0x80, 0x96)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Syne"
BODY_FONT = "DM Sans"


def rgb(hex_value):
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = CREAM
    stripe_h = Inches(0.09)
    for i, color in enumerate([EMERALD, RED, YELLOW]):
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(7.5 - 0.27 + i * 0.09),
            Inches(13.333),
            stripe_h,
        )
        set_fill(shp, color)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, font=BODY_FONT,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.65, 0.46, 8.8, 0.52, 24, FOREST, True, TITLE_FONT)
    if subtitle:
        add_text(slide, subtitle, 0.68, 0.99, 8.7, 0.32, 9.5, MUTED, False, BODY_FONT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.35), Inches(1.05), Inches(0.055))
    set_fill(line, YELLOW)


def add_footer(slide, idx, total, timing):
    add_text(slide, f"{idx:02d}/{total:02d}", 11.95, 7.08, 0.55, 0.22, 7.5, MUTED, False, BODY_FONT, PP_ALIGN.RIGHT)
    add_text(slide, timing, 10.5, 7.08, 1.25, 0.22, 7.5, MUTED, False, BODY_FONT, PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=16, color=INK, bullet_color=EMERALD, gap=0.43):
    for i, item in enumerate(items):
        cy = y + i * gap
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(cy + 0.08), Inches(0.10), Inches(0.10))
        set_fill(dot, bullet_color)
        add_text(slide, item, x + 0.22, cy, w - 0.22, 0.32, size, color, False, BODY_FONT)


def add_card(slide, x, y, w, h, title, body, accent=EMERALD):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CREAM_2
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(bar, accent)
    add_text(slide, title, x + 0.25, y + 0.18, w - 0.45, 0.28, 12.5, FOREST, True, TITLE_FONT)
    add_text(slide, body, x + 0.25, y + 0.58, w - 0.45, h - 0.70, 10.2, INK, False, BODY_FONT)


def add_metric(slide, x, y, number, label, color=FOREST):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.25), Inches(1.25))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = CREAM_2
    add_text(slide, number, x + 0.20, y + 0.15, 1.85, 0.45, 28, color, True, TITLE_FONT, PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.18, y + 0.76, 1.90, 0.30, 9.8, MUTED, True, BODY_FONT, PP_ALIGN.CENTER)


def add_notes(slide, text):
    try:
        notes = slide.notes_slide.notes_text_frame
        notes.text = text
    except Exception:
        pass


def add_logo(slide, x=10.25, y=0.35, w=2.2):
    candidates = [
        ROOT / "frontend/public/images/SafeTrip_4k_transparent.png",
        ROOT / "frontend/public/images/logo-removebg-preview (2).png",
        ROOT / "frontend/public/images/SafeTrip_4k_white_transparent.png",
    ]
    for path in candidates:
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
            return


slides = [
    {
        "title": "SafeTrip — Audit & Plan d'amélioration",
        "timing": "1 min",
        "notes": "Introduire l'objectif : transformer l'audit en plan de stabilisation concret. Le message clé est que SafeTrip a une base forte, mais nécessite des correctifs ciblés avant la bêta.",
    },
    {
        "title": "Agenda",
        "timing": "1 min",
        "bullets": ["Contexte produit et stack", "Résumé exécutif de l'audit", "Risques critiques à corriger", "Plan d'amélioration par horizon", "Décisions et prochaines actions"],
        "notes": "Présenter le déroulé en cinq blocs. Insister sur le fait que la réunion doit déboucher sur des assignations et une date de démo.",
    },
    {
        "title": "Contexte projet",
        "timing": "1 min",
        "notes": "Rappeler la promesse SafeTrip : réserver et suivre un voyage interurbain de façon plus fiable. Le parcours cible est volontairement simple : rechercher, réserver, confirmer.",
    },
    {
        "title": "Résumé exécutif",
        "timing": "1 min 30",
        "notes": "Mettre l'accent sur les chiffres : 51 points, dont 5 critiques. La conclusion n'est pas de ralentir le produit, mais de prioriser les corrections qui débloquent une bêta crédible.",
    },
    {
        "title": "Top 5 problèmes critiques",
        "timing": "2 min",
        "notes": "Lire cette slide par impacts, pas par détails de code. Certains points ont déjà évolué dans le code courant et doivent être vérifiés comme risques de régression.",
    },
    {
        "title": "Sécurité & risques",
        "timing": "1 min 30",
        "notes": "Séparer les actions immédiates des améliorations de maturité. Le fallback JWT et les routes sensibles sont les priorités de sécurité avant toute bêta externe.",
    },
    {
        "title": "UX/UI — points à améliorer",
        "timing": "1 min 15",
        "notes": "Le design possède déjà une identité forte. Le travail recommandé porte surtout sur la cohérence : tokens, boutons, feedbacks, chargements et validations.",
    },
    {
        "title": "Démo rapide — parcours utilisateur",
        "timing": "2 min",
        "notes": "Décrire le parcours comme un scénario utilisateur. Montrer où l'expérience doit être vérifiée : recherche homepage, sélection trajet, siège, paiement et confirmation.",
    },
    {
        "title": "Architecture & dépendances",
        "timing": "1 min 30",
        "notes": "Montrer que la dette principale vient du mélange API, fallback localStorage et configuration locale. La cible est une API source de vérité avec une configuration d'environnement nette.",
    },
    {
        "title": "Plan court terme — hotfix",
        "timing": "1 min 30",
        "notes": "Présenter ce plan comme le minimum pour une démo fiable. Chaque item doit avoir un propriétaire et être vérifié dans un parcours complet.",
    },
    {
        "title": "Plan moyen terme — bêta",
        "timing": "1 min 15",
        "notes": "La bêta doit valider l'usage réel, pas seulement l'affichage. Les tests, la persistance et la qualité mobile deviennent les critères de sortie.",
    },
    {
        "title": "Roadmap long terme — production",
        "timing": "1 min",
        "notes": "Mettre en avant la trajectoire de professionnalisation : paiement réel, monitoring, auth durcie et exploitation multi-agences.",
    },
    {
        "title": "Estimation effort & planning",
        "timing": "1 min 30",
        "notes": "Ce planning est une proposition de cadrage. Il doit être ajusté selon disponibilité équipe, accès aux credentials Supabase et niveau attendu de QA.",
    },
    {
        "title": "Dépendances & risques techniques",
        "timing": "1 min",
        "notes": "Cette slide prépare les arbitrages. Sans variables d'environnement, données fiables et tests, les mêmes bugs peuvent réapparaître à chaque démo.",
    },
    {
        "title": "Captures / assets requis",
        "timing": "45 sec",
        "notes": "Donner la liste des captures à fournir pour renforcer la présentation. Le deck est utilisable maintenant, mais les captures UI réelles amélioreront fortement la preuve visuelle.",
    },
    {
        "title": "Appel à l'action",
        "timing": "1 min",
        "notes": "Conclure par les décisions attendues : valider le périmètre hotfix, nommer les responsables et fixer la prochaine démo.",
    },
    {
        "title": "Annexe — Patchs critiques",
        "timing": "référence",
        "notes": "Cette annexe sert de feuille de route technique. Elle peut être utilisée directement pour créer les tickets de correction.",
    },
]


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = len(slides)

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = CREAM
    panel.line.color.rgb = CREAM
    add_logo(slide, 0.65, 0.55, 2.5)
    add_text(slide, "SafeTrip", 0.78, 2.0, 4.0, 0.55, 26, FOREST, True, TITLE_FONT)
    add_text(slide, "Audit & Plan d'amélioration", 0.78, 2.58, 7.6, 0.65, 34, FOREST, True, TITLE_FONT)
    add_text(slide, "Audit produit, sécurité et mise en production", 0.82, 3.45, 5.8, 0.35, 14, INK, False, BODY_FONT)
    add_text(slide, "27 mai 2026 · Équipe SafeTrip", 0.82, 4.05, 4.2, 0.28, 11, MUTED, True, BODY_FONT)
    for i, c in enumerate([EMERALD, RED, YELLOW]):
        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.5 + i * 0.25), Inches(0), Inches(0.18), Inches(7.5))
        set_fill(strip, c)
    add_card(slide, 8.9, 4.55, 3.45, 1.18, "Message clé", "Base produit prometteuse, mais stabilisation sécurité, configuration et UX nécessaire avant bêta.", YELLOW)
    add_footer(slide, 1, total, slides[0]["timing"])
    add_notes(slide, slides[0]["notes"])

    # 2 Agenda
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Agenda", "Cadrer rapidement les priorités pour passer de l'audit à l'exécution.")
    add_bullets(slide, slides[1]["bullets"], 0.95, 2.0, 6.8, 3.0, 19, FOREST, YELLOW, 0.62)
    add_card(slide, 8.1, 2.0, 3.9, 2.5, "Objectif réunion", "Valider le périmètre hotfix, assigner les responsables et fixer une date de démo bêta.", EMERALD)
    add_footer(slide, 2, total, slides[1]["timing"])
    add_notes(slide, slides[1]["notes"])

    # 3 Context
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Contexte projet", "Mission, stack et parcours critique à sécuriser.")
    add_card(slide, 0.8, 1.8, 3.6, 1.55, "Mission", "Simplifier la réservation de voyages interurbains et la traçabilité des bagages au Cameroun.", EMERALD)
    add_card(slide, 4.85, 1.8, 3.2, 1.55, "Stack", "Frontend Next.js 16 / React 19\nBackend Express TypeScript\nBase Supabase", YELLOW)
    add_card(slide, 8.5, 1.8, 3.2, 1.55, "Parcours cible", "Rechercher → Réserver → Confirmer\nLe flow doit être fiable de bout en bout.", RED)
    for i, label in enumerate(["Rechercher", "Réserver", "Confirmer"]):
        x = 1.15 + i * 3.8
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.55), Inches(2.55), Inches(0.68))
        pill.fill.solid()
        pill.fill.fore_color.rgb = FOREST if i == 0 else WHITE
        pill.line.color.rgb = FOREST
        add_text(slide, label, x, 4.73, 2.55, 0.22, 13, WHITE if i == 0 else FOREST, True, BODY_FONT, PP_ALIGN.CENTER)
        if i < 2:
            add_text(slide, "→", x + 2.85, 4.66, 0.35, 0.3, 20, EMERALD, True, TITLE_FONT)
    add_footer(slide, 3, total, slides[2]["timing"])
    add_notes(slide, slides[2]["notes"])

    # 4 Executive summary
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Résumé exécutif", "L'audit identifie une base exploitable, avec des risques à corriger avant bêta.")
    add_metric(slide, 0.9, 1.9, "51", "points à traiter", FOREST)
    add_metric(slide, 3.5, 1.9, "5", "bugs critiques", RED)
    add_metric(slide, 6.1, 1.9, "23", "fonctionnalités manquantes", EMERALD)
    add_metric(slide, 8.7, 1.9, "6", "risques sécurité", YELLOW)
    add_card(slide, 1.0, 4.15, 10.7, 1.35, "Diagnostic", "SafeTrip est suffisamment avancé pour une démonstration produit, mais les risques sécurité, configuration d'environnement et validation UX doivent être traités avant une bêta externe.", EMERALD)
    add_footer(slide, 4, total, slides[3]["timing"])
    add_notes(slide, slides[3]["notes"])

    # 5 Top critical
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Top 5 problèmes critiques", "Lecture par impact produit et risque de démonstration.")
    issues = [
        ("localStorage", "Collision / divergence cache, tickets et colis", "Données incohérentes"),
        ("PDF colis", "Interpolation à vérifier dans l'étiquette bagage", "Document inutilisable si régression"),
        ("Routes API", "Protection JWT incomplète sur routes sensibles", "Données exposées"),
        ("Recherche", "Paramètres homepage non transmis à la réservation", "Parcours interrompu"),
        ("Paiement", "Champs carte MM/AA et CVV non contrôlés", "Validation faible"),
    ]
    y = 1.75
    for label, body, impact in issues:
        add_card(slide, 0.85, y, 3.1, 0.68, label, body, RED if label in ["Routes API", "PDF colis"] else EMERALD)
        add_text(slide, impact, 4.25, y + 0.18, 4.0, 0.25, 12, FOREST, True, BODY_FONT)
        y += 0.86
    add_text(slide, "Note : vérifier les divergences entre le rapport initial et le code actuel avant création des tickets.", 0.9, 6.25, 9.2, 0.3, 10, MUTED, False, BODY_FONT)
    add_footer(slide, 5, total, slides[4]["timing"])
    add_notes(slide, slides[4]["notes"])

    # 6 Security
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Sécurité & risques", "Actions immédiates recommandées avant exposition bêta.")
    add_bullets(slide, [
        "Rendre JWT_SECRET obligatoire : aucun fallback en production.",
        "Appliquer requireAuth / requireRole sur les routes agence, client et admin.",
        "Restreindre CORS aux domaines autorisés.",
        "Réduire les informations exposées par /health.",
        "Supprimer les accès de simulation admin hors environnement local.",
    ], 0.9, 1.8, 6.4, 3.2, 15, FOREST, RED, 0.52)
    add_card(slide, 8.0, 1.75, 3.75, 2.95, "Fichiers ciblés", "backend/src/controllers/authController.ts\nbackend/src/middleware/authMiddleware.ts\nbackend/src/index.ts\nbackend/src/routes/*Routes.ts", RED)
    add_footer(slide, 6, total, slides[5]["timing"])
    add_notes(slide, slides[5]["notes"])

    # 7 UX/UI
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "UX/UI — points à améliorer", "Conserver l'identité, systématiser l'expérience.")
    add_card(slide, 0.9, 1.75, 3.25, 1.45, "Design tokens", "Centraliser couleurs, spacing, radius, ombres et typographies dans une base commune.", EMERALD)
    add_card(slide, 4.45, 1.75, 3.25, 1.45, "Boutons & CTA", "Unifier états : normal, hover, focus, disabled, loading.", YELLOW)
    add_card(slide, 8.0, 1.75, 3.25, 1.45, "Formulaires", "Validations inline, focus rings, messages d'erreur lisibles.", RED)
    add_card(slide, 2.7, 4.0, 3.25, 1.45, "Loading", "Skeletons et états vides conçus pour éviter les ruptures de confiance.", EMERALD)
    add_card(slide, 6.25, 4.0, 3.25, 1.45, "Mobile", "Réduire la charge cognitive et optimiser dashboards et recherche.", YELLOW)
    add_footer(slide, 7, total, slides[6]["timing"])
    add_notes(slide, slides[6]["notes"])

    # 8 Demo
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Démo rapide — parcours utilisateur", "Zones à vérifier pendant la présentation.")
    steps = [
        ("1", "Homepage", "Hero + recherche : départ, destination, date"),
        ("2", "Résultats", "Cartes trajets, filtres, CTA Réserver"),
        ("3", "Siège", "Plan, disponibilité et sélection"),
        ("4", "Paiement", "OM / MoMo / Carte, validation champs"),
        ("5", "Confirmation", "Billet, historique, colis, PDF"),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.8 + i * 2.45
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.55), Inches(1.85), Inches(0.65), Inches(0.65))
        set_fill(circ, FOREST if i == 0 else WHITE)
        circ.line.color.rgb = FOREST
        add_text(slide, num, x + 0.55, 2.02, 0.65, 0.2, 12, WHITE if i == 0 else FOREST, True, BODY_FONT, PP_ALIGN.CENTER)
        add_text(slide, title, x, 2.75, 1.75, 0.26, 12.5, FOREST, True, TITLE_FONT, PP_ALIGN.CENTER)
        add_text(slide, body, x - 0.10, 3.18, 1.95, 0.55, 8.6, MUTED, False, BODY_FONT, PP_ALIGN.CENTER)
    add_card(slide, 1.3, 5.0, 10.2, 0.85, "Fichiers", "frontend/src/app/page.tsx · frontend/src/app/reserver/page.tsx · frontend/src/app/reserver/confirmer/page.tsx · frontend/src/app/client/dashboard/page.tsx", EMERALD)
    add_footer(slide, 8, total, slides[7]["timing"])
    add_notes(slide, slides[7]["notes"])

    # 9 Architecture
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Architecture & dépendances", "Cible : API source de vérité et configuration externalisée.")
    nodes = [
        (0.9, 2.3, "Next.js", "UI / parcours réservation"),
        (4.1, 2.3, "Express API", "Auth, agences, clients"),
        (7.3, 2.3, "Supabase", "Données persistantes"),
        (4.1, 4.35, "LocalStorage", "Fallback temporaire"),
    ]
    for x, y, title, body in nodes:
        add_card(slide, x, y, 2.55, 1.05, title, body, EMERALD if title != "LocalStorage" else YELLOW)
    add_text(slide, "→", 3.63, 2.65, 0.25, 0.25, 20, FOREST, True, TITLE_FONT)
    add_text(slide, "→", 6.83, 2.65, 0.25, 0.25, 20, FOREST, True, TITLE_FONT)
    add_text(slide, "Fallback", 4.95, 3.78, 1.2, 0.25, 10, MUTED, True, BODY_FONT, PP_ALIGN.CENTER)
    add_card(slide, 9.95, 2.25, 2.2, 2.2, "À corriger", "Centraliser API_BASE\nSupprimer fallback prod\nCréer .env examples\nVersionner l'API", RED)
    add_footer(slide, 9, total, slides[8]["timing"])
    add_notes(slide, slides[8]["notes"])

    # 10 Short term
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Plan court terme — hotfix", "Objectif : démo stable et parcours critique fiable.")
    add_bullets(slide, [
        "Supprimer fallback JWT et configurer les secrets.",
        "Protéger les routes sensibles avec JWT + rôles.",
        "Contrôler les champs paiement carte et mobile money.",
        "Vérifier PDF colis et nom passager dynamique.",
        "Harmoniser localStorage et API.",
        "Créer .env.local.example et backend/.env.example.",
    ], 0.9, 1.75, 6.8, 3.5, 14.5, FOREST, RED, 0.48)
    add_metric(slide, 8.35, 2.0, "S1", "semaine hotfix", RED)
    add_metric(slide, 8.35, 3.55, "6", "patchs prioritaires", FOREST)
    add_footer(slide, 10, total, slides[9]["timing"])
    add_notes(slide, slides[9]["notes"])

    # 11 Mid term
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Plan moyen terme — bêta", "Objectif : valider l'usage réel sans casser le parcours.")
    add_bullets(slide, [
        "Tests E2E du parcours Réserver → Confirmer.",
        "API comme source de vérité pour tickets et colis.",
        "Dashboard client fiable avec données synchronisées.",
        "Design system appliqué aux composants critiques.",
        "Observabilité minimale : logs, erreurs et santé API.",
    ], 0.95, 1.85, 7.0, 3.3, 15.2, FOREST, EMERALD, 0.55)
    add_card(slide, 8.3, 2.15, 3.35, 1.85, "Critères de sortie", "Réservation complète testée\nAucun secret en fallback\nDonnées persistées\nDémo mobile validée", YELLOW)
    add_footer(slide, 11, total, slides[10]["timing"])
    add_notes(slide, slides[10]["notes"])

    # 12 Long term
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Roadmap long terme — production", "Industrialiser le produit et préparer l'exploitation.")
    add_card(slide, 0.9, 1.85, 3.0, 1.4, "Paiement réel", "Passerelle CinetPay / PayDunya, rapprochement transactionnel.", EMERALD)
    add_card(slide, 4.15, 1.85, 3.0, 1.4, "Traçabilité", "QR codes scannables, scanner agent, statuts colis fiables.", YELLOW)
    add_card(slide, 7.4, 1.85, 3.0, 1.4, "Production", "CI/CD, monitoring, sauvegardes, alerting et durcissement auth.", RED)
    add_card(slide, 2.5, 4.1, 3.0, 1.4, "Multi-agences", "Gestion opérations, modification trajets, suppression bus.", EMERALD)
    add_card(slide, 5.75, 4.1, 3.0, 1.4, "Client", "Notifications, email, annulation, remboursement et profil sync.", YELLOW)
    add_footer(slide, 12, total, slides[11]["timing"])
    add_notes(slide, slides[11]["notes"])

    # 13 Planning
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Estimation effort & planning", "Proposition de séquencement pour 5 semaines.")
    stages = [
        ("S1", "Hotfix critiques", RED),
        ("S2", "Bêta interne", EMERALD),
        ("S3-S4", "Bêta externe", YELLOW),
        ("S5+", "Durcissement prod", FOREST),
    ]
    for i, (week, label, color) in enumerate(stages):
        x = 0.95 + i * 2.95
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.5), Inches(2.35), Inches(1.1))
        set_fill(rect, color)
        add_text(slide, week, x + 0.15, 2.72, 2.0, 0.28, 16, WHITE, True, TITLE_FONT, PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.15, 3.18, 2.0, 0.25, 9.5, WHITE, True, BODY_FONT, PP_ALIGN.CENTER)
    add_card(slide, 1.25, 4.75, 10.45, 0.85, "Hypothèse", "Planning à ajuster selon disponibilité équipe, accès Supabase, niveau de test attendu et validation métier.", EMERALD)
    add_footer(slide, 13, total, slides[12]["timing"])
    add_notes(slide, slides[12]["notes"])

    # 14 Risks
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Dépendances & risques techniques", "Points qui peuvent ralentir ou fragiliser la bêta.")
    add_bullets(slide, [
        "Credentials Supabase et variables d'environnement non stabilisés.",
        "Données locales et API potentiellement divergentes.",
        "Absence de tests automatisés sur le flow réservation.",
        "Routes sensibles encore à protéger et valider.",
        "Risque de régression UI sans design system partagé.",
    ], 0.9, 1.85, 7.2, 3.2, 15, FOREST, YELLOW, 0.56)
    add_card(slide, 8.4, 2.0, 3.15, 2.35, "Réduction du risque", "Créer tickets atomiques\nAjouter smoke tests\nMettre une checklist PR\nValider en démo quotidienne", EMERALD)
    add_footer(slide, 14, total, slides[13]["timing"])
    add_notes(slide, slides[13]["notes"])

    # 15 Assets
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Captures / assets requis", "Images à ajouter pour renforcer la preuve visuelle.")
    add_bullets(slide, [
        "01_home_hero.png — hero + capsule recherche.",
        "02_reserver_list.png — trajets, filtres, bouton Réserver.",
        "03_confirmer_seats.png — plan des sièges.",
        "04_payment.png — paiement OM / MoMo / carte.",
        "05_pdf_label_check.png — étiquette colis PDF.",
        "07_backend_jwt.png et 09_api_url_occurrences.png — extraits sécurité/config.",
    ], 0.9, 1.75, 7.9, 3.8, 12.8, FOREST, EMERALD, 0.45)
    add_card(slide, 9.0, 2.05, 2.75, 2.5, "Dossier", "outputs/screenshots/\nPNG 1920×1080 pour UI\nPNG 1200×800 pour code", YELLOW)
    add_footer(slide, 15, total, slides[14]["timing"])
    add_notes(slide, slides[14]["notes"])

    # 16 CTA
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Appel à l'action", "Décisions attendues aujourd'hui.")
    add_card(slide, 0.95, 1.9, 3.1, 1.55, "Valider", "Périmètre des 6 patchs critiques et critères de sortie.", EMERALD)
    add_card(slide, 4.35, 1.9, 3.1, 1.55, "Assigner", "Responsables frontend, backend et QA.", YELLOW)
    add_card(slide, 7.75, 1.9, 3.1, 1.55, "Planifier", "Date de démo bêta et revue technique.", RED)
    add_text(slide, "Prochaine démo cible : à confirmer", 1.15, 4.85, 7.2, 0.5, 22, FOREST, True, TITLE_FONT)
    add_text(slide, "Décision recommandée : lancer le sprint hotfix avant toute bêta externe.", 1.18, 5.45, 7.5, 0.3, 12, MUTED, False, BODY_FONT)
    add_footer(slide, 16, total, slides[15]["timing"])
    add_notes(slide, slides[15]["notes"])

    # 17 Appendix
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Annexe — Patchs critiques", "6 corrections directement transformables en tickets.")
    patches = [
        ("localStorage", "Harmoniser cache trajets, historique billets et colis.", "reserver/page.tsx · confirmer/page.tsx · client/dashboard/page.tsx"),
        ("PDF colis", "Vérifier les interpolations et retirer tout texte hardcodé.", "client/dashboard/page.tsx"),
        ("JWT", "Supprimer fallback secret et protéger routes sensibles.", "authController.ts · authMiddleware.ts · routes/*"),
        ("Paiement", "Contrôler MM/AA, CVV, carte, téléphone et référence.", "reserver/confirmer/page.tsx"),
        ("API URL", "Centraliser NEXT_PUBLIC_API_URL, retirer fallback prod.", "frontend/src/app/**/*.tsx"),
        ("Env", "Ajouter exemples d'environnement frontend et backend.", "frontend/.env.local.example · backend/.env.example"),
    ]
    for i, (title, action, files) in enumerate(patches):
        y = 1.6 + i * 0.75
        add_text(slide, f"{i + 1}", 0.85, y + 0.08, 0.32, 0.25, 11, WHITE, True, BODY_FONT, PP_ALIGN.CENTER)
        num_bg = slide.shapes[-1]
        # textboxes cannot be filled reliably; draw number chip behind afterward
        chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.78), Inches(y + 0.02), Inches(0.42), Inches(0.42))
        set_fill(chip, FOREST)
        add_text(slide, f"{i + 1}", 0.78, y + 0.11, 0.42, 0.16, 8.5, WHITE, True, BODY_FONT, PP_ALIGN.CENTER)
        add_text(slide, title, 1.35, y, 1.45, 0.25, 11, FOREST, True, TITLE_FONT)
        add_text(slide, action, 2.75, y, 4.4, 0.25, 9.8, INK, False, BODY_FONT)
        add_text(slide, files, 7.35, y, 4.7, 0.25, 8.6, MUTED, False, BODY_FONT)
    add_footer(slide, 17, total, slides[16]["timing"])
    add_notes(slide, slides[16]["notes"])

    prs.save(PPTX_PATH)


def build_pdf():
    c = canvas.Canvas(str(PDF_PATH), pagesize=A4)
    width, height = A4
    margin = 18 * mm
    y = height - margin

    c.setFillColor(colors.HexColor("#F7F4EE"))
    c.rect(0, 0, width, height, stroke=0, fill=1)
    c.setFillColor(colors.HexColor("#0A2F1D"))
    c.rect(0, height - 9 * mm, width, 9 * mm, stroke=0, fill=1)
    c.setFillColor(colors.HexColor("#CE1126"))
    c.rect(0, height - 12 * mm, width, 3 * mm, stroke=0, fill=1)
    c.setFillColor(colors.HexColor("#FCD116"))
    c.rect(0, height - 15 * mm, width, 3 * mm, stroke=0, fill=1)

    c.setFillColor(colors.HexColor("#0A2F1D"))
    c.setFont("Helvetica-Bold", 20)
    c.drawString(margin, y - 14 * mm, "Checklist technique — SafeTrip")
    c.setFont("Helvetica", 9)
    c.setFillColor(colors.HexColor("#718096"))
    c.drawString(margin, y - 20 * mm, "Audit & Plan d'amélioration · 27 mai 2026")

    items = [
        "Harmoniser les clés localStorage : cache trajets vs historique billets.",
        "Vérifier/corriger l'interpolation dans l'étiquette bagage PDF.",
        "Supprimer le fallback JWT secret et rendre JWT_SECRET obligatoire.",
        "Renforcer le bouton Réserver et le parcours de sélection trajet.",
        "Contrôler et valider tous les inputs paiement : MM/AA, CVV, carte, téléphone.",
        "Remplacer les URLs http://localhost:5000 par NEXT_PUBLIC_API_URL ou une config centralisée.",
        "Ajouter frontend/.env.local.example et backend/.env.example.",
        "Restreindre CORS en production.",
        "Protéger les routes API sensibles avec requireAuth / requireRole.",
        "Nettoyer /health pour ne pas exposer d'informations sensibles.",
    ]

    y = height - 55 * mm
    c.setFont("Helvetica", 10.5)
    c.setFillColor(colors.HexColor("#1C2B22"))
    for item in items:
        c.setStrokeColor(colors.HexColor("#0A2F1D"))
        c.setLineWidth(1)
        c.rect(margin, y - 2 * mm, 4 * mm, 4 * mm, stroke=1, fill=0)
        c.drawString(margin + 8 * mm, y - 1 * mm, item)
        y -= 12 * mm

    c.setFillColor(colors.HexColor("#EDE9DF"))
    c.roundRect(margin, 24 * mm, width - 2 * margin, 22 * mm, 4 * mm, stroke=0, fill=1)
    c.setFillColor(colors.HexColor("#0A2F1D"))
    c.setFont("Helvetica-Bold", 10)
    c.drawString(margin + 6 * mm, 38 * mm, "Validation recommandée")
    c.setFont("Helvetica", 9)
    c.drawString(margin + 6 * mm, 32 * mm, "Créer un ticket par correction, assigner un propriétaire et vérifier le parcours complet en démo.")

    c.save()


def add_table(slide, x, y, widths, row_h, headers, rows, accent=FOREST, font_size=8.5):
    cur_x = x
    for i, head in enumerate(headers):
        rect = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(cur_x),
            Inches(y),
            Inches(widths[i]),
            Inches(row_h),
        )
        set_fill(rect, accent)
        add_text(slide, head, cur_x + 0.08, y + 0.12, widths[i] - 0.16, row_h - 0.12, 8.5, WHITE, True, BODY_FONT)
        cur_x += widths[i]

    for r, row in enumerate(rows):
        cur_x = x
        bg_color = WHITE if r % 2 == 0 else CREAM_2
        for i, value in enumerate(row):
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(cur_x),
                Inches(y + row_h * (r + 1)),
                Inches(widths[i]),
                Inches(row_h),
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = bg_color
            rect.line.color.rgb = CREAM_2
            add_text(slide, value, cur_x + 0.08, y + row_h * (r + 1) + 0.09, widths[i] - 0.16, row_h - 0.1, font_size, INK, i == 0, BODY_FONT)
            cur_x += widths[i]


def build_pptx_10():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 10

    compact_slides = [
        ("Couverture", "1 min", "Introduire l'objectif : transformer l'audit en plan de stabilisation concret avant bêta."),
        ("Agenda & contexte", "1 min 30", "Présenter la mission SafeTrip, la stack et le fil conducteur de la présentation."),
        ("Résumé exécutif", "1 min 30", "Faire comprendre rapidement l'ampleur de l'audit et le message produit principal."),
        ("Top 5 problèmes critiques", "2 min", "Lire par impact et priorité. Les divergences entre rapport et code doivent rester marquées comme à vérifier."),
        ("Sécurité & architecture", "2 min", "Relier les risques sécurité aux dépendances techniques : JWT, CORS, health endpoint, API URL et localStorage."),
        ("UX/UI & démo parcours", "2 min", "Raconter le parcours utilisateur et les zones à valider pendant une démo."),
        ("Plan court terme", "1 min 30", "Présenter les hotfix comme minimum nécessaire pour une bêta testable."),
        ("Bêta & production", "1 min 30", "Montrer la progression entre stabilisation bêta et industrialisation production."),
        ("Planning & risques", "1 min 30", "Clarifier le calendrier cible, les estimations et les dépendances critiques."),
        ("Action + patchs critiques", "2 min", "Conclure avec les décisions attendues et la liste des tickets techniques à créer."),
    ]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_logo(slide, 0.65, 0.55, 2.5)
    add_text(slide, "SafeTrip", 0.78, 1.9, 4.0, 0.55, 26, FOREST, True, TITLE_FONT)
    add_text(slide, "Audit & Plan d'amélioration", 0.78, 2.52, 8.4, 0.68, 34, FOREST, True, TITLE_FONT)
    add_text(slide, "Audit produit, sécurité et mise en production", 0.82, 3.42, 5.8, 0.35, 14, INK, False, BODY_FONT)
    add_text(slide, "27 mai 2026 · Équipe SafeTrip", 0.82, 4.02, 4.2, 0.28, 11, MUTED, True, BODY_FONT)
    for i, c in enumerate([EMERALD, RED, YELLOW]):
        strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.65 + i * 0.25), Inches(0), Inches(0.18), Inches(7.5))
        set_fill(strip, c)
    add_card(slide, 9.05, 4.55, 3.15, 1.18, "Message clé", "Produit prometteur, mais stabilisation sécurité, configuration et UX nécessaire avant bêta.", YELLOW)
    add_footer(slide, 1, total, compact_slides[0][1])
    add_notes(slide, compact_slides[0][2])

    # 2 Agenda & context
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Agenda & contexte", "Une lecture orientée décision : risques, plan, responsabilités.")
    add_bullets(slide, ["Contexte", "Résumé exécutif", "Risques critiques", "Plan d'amélioration", "Décisions attendues"], 0.9, 1.72, 4.6, 3.2, 16, FOREST, YELLOW, 0.52)
    add_card(slide, 5.45, 1.75, 2.95, 1.35, "Mission", "Réservation de voyages interurbains et traçabilité bagages au Cameroun.", EMERALD)
    add_card(slide, 8.75, 1.75, 2.95, 1.35, "Stack", "Next.js 16 / React 19\nExpress TypeScript\nSupabase", YELLOW)
    for i, label in enumerate(["Rechercher", "Réserver", "Confirmer"]):
        x = 4.15 + i * 2.45
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.65), Inches(1.85), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = FOREST if i == 0 else WHITE
        pill.line.color.rgb = FOREST
        add_text(slide, label, x, 4.80, 1.85, 0.18, 9.5, WHITE if i == 0 else FOREST, True, BODY_FONT, PP_ALIGN.CENTER)
        if i < 2:
            add_text(slide, "→", x + 2.05, 4.78, 0.25, 0.2, 15, EMERALD, True, TITLE_FONT)
    add_card(slide, 1.0, 5.75, 10.8, 0.6, "Capture à intégrer si disponible", "outputs/screenshots/01_home_hero.png — hero + capsule de recherche.", EMERALD)
    add_footer(slide, 2, total, compact_slides[1][1])
    add_notes(slide, compact_slides[1][2])

    # 3 Exec summary
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Résumé exécutif", "51 points relevés, 5 critiques, 23 fonctionnalités manquantes.")
    add_metric(slide, 0.9, 1.8, "51", "points relevés", FOREST)
    add_metric(slide, 3.55, 1.8, "5", "critiques", RED)
    add_metric(slide, 6.2, 1.8, "23", "manquantes", EMERALD)
    add_metric(slide, 8.85, 1.8, "6", "risques sécurité", YELLOW)
    add_card(slide, 1.0, 4.05, 10.7, 1.3, "Message principal", "SafeTrip est démontrable et dispose d'une identité produit solide. Avant bêta, il faut sécuriser les secrets, clarifier la configuration, fiabiliser le parcours réservation et réduire les divergences localStorage/API.", EMERALD)
    add_footer(slide, 3, total, compact_slides[2][1])
    add_notes(slide, compact_slides[2][2])

    # 4 Top 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Top 5 problèmes critiques", "Tableau d'impact et de priorité.")
    add_table(
        slide,
        0.75,
        1.6,
        [2.25, 6.0, 2.25, 1.25],
        0.55,
        ["Problème", "Impact", "Priorité", "Statut"],
        [
            ["JWT fallback", "Secret hardcodé si .env absent ; tokens prévisibles.", "P0", "Actif"],
            ["Routes sensibles", "Protection JWT/rôle à appliquer partout.", "P0", "À vérifier"],
            ["Paiement carte", "MM/AA et CVV visibles mais non contrôlés.", "P1", "Actif"],
            ["localStorage/API", "Risque de divergence cache, billets, colis.", "P1", "À vérifier"],
            ["Recherche homepage", "Les paramètres saisis ne sont pas transmis au flow réservation.", "P1", "Actif"],
        ],
        RED,
        8.2,
    )
    add_text(slide, "Règle : si le rapport et le code divergent, créer un ticket de vérification avant correction.", 0.82, 6.35, 9.3, 0.25, 9.5, MUTED, False, BODY_FONT)
    add_footer(slide, 4, total, compact_slides[3][1])
    add_notes(slide, compact_slides[3][2])

    # 5 Security architecture
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Sécurité & architecture", "Risques immédiats et dépendances de déploiement.")
    add_card(slide, 0.85, 1.65, 3.25, 1.45, "Sécurité", "JWT_SECRET obligatoire\nMiddleware existant mais fallback risqué\nCORS ouvert via app.use(cors())", RED)
    add_card(slide, 4.45, 1.65, 3.25, 1.45, "Architecture", "Next.js → Express API → Supabase\nLocalStorage comme fallback temporaire", EMERALD)
    add_card(slide, 8.05, 1.65, 3.25, 1.45, "Configuration", "NEXT_PUBLIC_API_URL || http://localhost:5000\n.env.local et backend/.env à cadrer", YELLOW)
    add_bullets(slide, [
        "backend/src/controllers/authController.ts",
        "backend/src/middleware/authMiddleware.ts",
        "backend/src/index.ts",
        "frontend/src/app/login|reserver|confirmer|client|agence|admin|agences|tracabilite/page.tsx",
    ], 1.05, 4.0, 10.2, 1.9, 11.5, FOREST, EMERALD, 0.42)
    add_footer(slide, 5, total, compact_slides[4][1])
    add_notes(slide, compact_slides[4][2])

    # 6 UX & demo
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "UX/UI & démo parcours", "Ce qu'il faut montrer et améliorer.")
    ux_items = [
        ("Tokens", "Couleurs, spacing, radius, ombres, typographie."),
        ("Loading", "Skeletons, états vides, feedbacks de chargement."),
        ("Boutons", "CTA cohérents, disabled lisible, focus visible."),
        ("Mobile", "Dashboards moins lourds, navigation plus efficace."),
        ("Formulaires", "Validation inline et accessibilité AA."),
    ]
    for i, (t, b) in enumerate(ux_items):
        add_card(slide, 0.8 + (i % 3) * 3.55, 1.65 + (i // 3) * 1.45, 3.1, 1.05, t, b, [EMERALD, YELLOW, RED, EMERALD, YELLOW][i])
    add_card(slide, 1.15, 5.0, 10.3, 0.85, "Démo à valider", "Recherche homepage → liste trajets → sélection trajet → siège → paiement → confirmation.", EMERALD)
    add_footer(slide, 6, total, compact_slides[5][1])
    add_notes(slide, compact_slides[5][2])

    # 7 Short term
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Plan court terme — hotfix prioritaires", "Objectif : rendre la bêta testable.")
    add_table(
        slide,
        0.8,
        1.65,
        [3.0, 4.8, 2.0, 1.15],
        0.52,
        ["Chantier", "Action", "Fichiers", "Effort"],
        [
            ["Secrets JWT", "Forcer JWT_SECRET, supprimer fallback.", "backend/*auth*", "S"],
            ["Variables env", "Créer exemples .env et centraliser API URL.", "frontend + backend", "S"],
            ["Paiement", "Valider carte, MM/AA, CVV, téléphone, référence.", "confirmer/page.tsx", "M"],
            ["PDF bagage", "Vérifier interpolation et données passager.", "client/dashboard", "S"],
            ["Routes sensibles", "requireAuth / requireRole sur API.", "routes/*.ts", "M"],
            ["localStorage/API", "Clarifier cache vs historique vs colis.", "3 pages frontend", "M"],
        ],
        EMERALD,
        7.4,
    )
    add_footer(slide, 7, total, compact_slides[6][1])
    add_notes(slide, compact_slides[6][2])

    # 8 Beta & production
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Bêta & production", "Deux horizons : stabiliser puis industrialiser.")
    add_card(slide, 0.95, 1.8, 4.9, 2.3, "Moyen terme — bêta", "Tests E2E réservation\nDashboard client fiable\nAPI source de vérité\nObservabilité minimale\nDesign system appliqué", EMERALD)
    add_card(slide, 6.55, 1.8, 4.9, 2.3, "Long terme — production", "Auth sécurisée\nMonitoring + CI/CD\nMigration Supabase complète\nPerformance mobile\nSupport multi-agences", FOREST)
    add_card(slide, 2.1, 5.0, 8.3, 0.75, "Critère commun", "Chaque release doit préserver le parcours Rechercher → Réserver → Confirmer.", YELLOW)
    add_footer(slide, 8, total, compact_slides[7][1])
    add_notes(slide, compact_slides[7][2])

    # 9 Planning & risks
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Planning & risques", "Timeline cible et dépendances à lever.")
    stages = [("S1", "Hotfix critiques", RED), ("S2", "Bêta interne", EMERALD), ("S3-S4", "Bêta externe", YELLOW), ("S5+", "Durcissement prod", FOREST)]
    for i, (week, label, color) in enumerate(stages):
        x = 0.85 + i * 2.9
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.75), Inches(2.35), Inches(1.0))
        set_fill(rect, color)
        add_text(slide, week, x + 0.15, 1.93, 2.0, 0.25, 14, WHITE, True, TITLE_FONT, PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.15, 2.35, 2.0, 0.22, 8.8, WHITE, True, BODY_FONT, PP_ALIGN.CENTER)
    add_bullets(slide, [
        "Variables d'environnement manquantes.",
        "Qualité et couverture des données Supabase.",
        "Dépendance localStorage pendant la transition.",
        "Absence de tests automatisés du flow critique.",
        "Risque de régression UI sans tokens partagés.",
    ], 1.05, 4.0, 7.4, 2.1, 12.5, FOREST, RED, 0.42)
    add_card(slide, 8.6, 4.1, 2.8, 1.45, "Mitigation", "Tickets atomiques\nChecklist PR\nSmoke test démo\nRevue quotidienne", EMERALD)
    add_footer(slide, 9, total, compact_slides[8][1])
    add_notes(slide, compact_slides[8][2])

    # 10 CTA appendix
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "Action + patchs critiques", "Décisions attendues et tickets à créer.")
    add_card(slide, 0.85, 1.55, 3.05, 0.9, "Valider", "Les 6 patchs critiques.", EMERALD)
    add_card(slide, 4.15, 1.55, 3.05, 0.9, "Assigner", "Frontend / Backend / QA.", YELLOW)
    add_card(slide, 7.45, 1.55, 3.05, 0.9, "Planifier", "Date de démo bêta.", RED)
    add_table(
        slide,
        0.75,
        3.0,
        [2.1, 6.4, 2.85],
        0.42,
        ["Patch", "Action précise", "Fichiers"],
        [
            ["localStorage", "Harmoniser cache trajets vs historique billets.", "reserver / confirmer / client"],
            ["PDF bagage", "Corriger toute interpolation échappée.", "client/dashboard"],
            ["JWT", "Supprimer fallback et forcer JWT_SECRET.", "authController / middleware"],
            ["Paiement", "Valider MM/AA, CVV, carte, téléphone, ref.", "confirmer/page.tsx"],
            ["API URL", "Remplacer fallback localhost par config centralisée.", "frontend app"],
            ["Env", "Ajouter .env.local.example et .env.example.", "frontend / backend"],
        ],
        FOREST,
        6.5,
    )
    add_text(slide, "Captures : 01_home_hero.png à 10_uiux_tokens.png dans outputs/screenshots/.", 0.85, 6.78, 9.5, 0.22, 8.5, MUTED, False, BODY_FONT)
    add_footer(slide, 10, total, compact_slides[9][1])
    add_notes(slide, compact_slides[9][2])

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    build_pptx_10()
    build_pdf()
    print(PPTX_PATH)
    print(PDF_PATH)
